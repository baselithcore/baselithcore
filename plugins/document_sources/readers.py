"""
Document Readers.

Provides specialized readers for different file formats (PDF, Markdown, etc.).
"""

from __future__ import annotations

from core.observability.logging import get_logger
from pathlib import Path

from typing import List, Optional

from .ocr_backends import run_image_ocr, run_pdf_ocr
from .utils import normalize_text, strip_front_matter, warn_missing_dependency

logger = get_logger(__name__)


def read_markdown(path: Path) -> Optional[str]:
    """
    Legge file Markdown rimuovendo eventuale front matter.
    Blocking I/O: This function performs synchronous file I/O.
    """

    try:
        raw = path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        try:
            raw = path.read_text(encoding="utf-8", errors="ignore")
        except OSError as exc:
            logger.warning(f"[filesystem] Error reading {path}: {exc}")
            return None
    except OSError as exc:
        logger.warning(f"[filesystem] Error reading {path}: {exc}")
        return None

    normalized = strip_front_matter(raw)
    cleaned = normalize_text(normalized)
    return cleaned or None


def read_pdf(path: Path) -> Optional[str]:
    """
    Legge testo da PDF, se serve tenta OCR.
    Blocking I/O: This function performs synchronous file I/O and CPU-bound parsing.
    """

    try:
        from pypdf import PdfReader
    except ImportError:  # pragma: no cover - dipendenza opzionale
        warn_missing_dependency("pypdf", "lettura testi PDF")
        return None

    text_parts: List[str] = []
    try:
        reader = PdfReader(str(path))
        for page in reader.pages:
            snippet = page.extract_text() or ""
            snippet = normalize_text(snippet)
            if snippet:
                text_parts.append(snippet)
    except Exception as exc:
        logger.warning(f"[filesystem] Error reading PDF {path}: {exc}")
        return None

    combined = "\n\n".join(part for part in text_parts if part)
    if combined.strip():
        return combined

    return run_pdf_ocr(path)


def read_image(path: Path) -> Optional[str]:
    """
    Legge testo da immagini tramite OCR.
    Blocking I/O: This function performs synchronous file I/O and CPU-bound OCR.
    """
    return run_image_ocr(path)


def read_word(path: Path) -> Optional[str]:
    """
    Legge documenti Word (docx/doc).
    Blocking I/O: This function performs synchronous file I/O.
    """

    suffix = path.suffix.lower()
    if suffix == ".docx":
        return _read_docx(path)
    if suffix == ".doc":
        return _read_doc_legacy(path)
    return None


def _read_docx(path: Path) -> Optional[str]:
    try:
        from docx import Document  # type: ignore[import]
    except ImportError:  # pragma: no cover - dipendenza opzionale
        warn_missing_dependency("python-docx", "lettura documenti Word")
        return None

    try:
        document = Document(str(path))
    except Exception as exc:
        logger.warning(f"[filesystem] Error reading Word {path}: {exc}")
        return None

    text_parts: List[str] = []
    for paragraph in document.paragraphs:
        snippet = normalize_text(paragraph.text)
        if snippet:
            text_parts.append(snippet)

    for table in document.tables:
        for row in table.rows:
            cells = [normalize_text(cell.text) for cell in row.cells if cell.text]
            row_text = " | ".join(filter(None, cells))
            if row_text:
                text_parts.append(row_text)

    combined = "\n\n".join(text_parts)
    return combined.strip() or None


def _read_doc_legacy(path: Path) -> Optional[str]:
    try:
        import textract  # type: ignore[import]
    except ImportError:  # pragma: no cover - dipendenza opzionale
        warn_missing_dependency("textract", "lettura documenti Word .doc")
        return None

    try:
        raw_bytes = textract.process(str(path))
    except Exception as exc:
        logger.warning(f"[filesystem] Error reading Word {path}: {exc}")
        return None

    try:
        text = raw_bytes.decode("utf-8")
    except UnicodeDecodeError:
        text = raw_bytes.decode("latin-1", errors="ignore")
    return normalize_text(text) or None


def read_powerpoint(path: Path) -> Optional[str]:
    """
    Estrae testo da presentazioni PowerPoint.
    Blocking I/O: This function performs synchronous file I/O.
    """

    try:
        from pptx import Presentation  # type: ignore[import]
    except ImportError:  # pragma: no cover - dipendenza opzionale
        warn_missing_dependency("python-pptx", "lettura presentazioni PowerPoint")
        return None

    try:
        presentation = Presentation(str(path))
    except Exception as exc:
        logger.warning(f"[filesystem] Error reading PowerPoint {path}: {exc}")
        return None

    text_parts: List[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_lines: List[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame:
                paragraphs: List[str] = []
                for paragraph in shape.text_frame.paragraphs:
                    snippet = normalize_text(paragraph.text)
                    if snippet:
                        paragraphs.append(snippet)
                if paragraphs:
                    slide_lines.append("\n".join(paragraphs))
                continue

            if getattr(shape, "has_table", False):
                table_lines: List[str] = []
                for row in shape.table.rows:  # type: ignore[union-attr]
                    cells: List[str] = []
                    for cell in row.cells:
                        snippet = normalize_text(cell.text)
                        if snippet:
                            cells.append(snippet)
                    if cells:
                        table_lines.append(" | ".join(cells))
                if table_lines:
                    slide_lines.append("\n".join(table_lines))
                continue

            text = getattr(shape, "text", "")
            snippet = normalize_text(text)
            if snippet:
                slide_lines.append(snippet)
        if slide_lines:
            slide_text = f"[Slide {index}]\n" + "\n".join(slide_lines)
            text_parts.append(slide_text)

    combined = "\n\n".join(text_parts)
    return combined.strip() or None


def read_excel(path: Path) -> Optional[str]:
    """
    Estrae dati leggibili da file Excel (xlsx/xls).
    Blocking I/O: This function performs synchronous file I/O.
    """

    suffix = path.suffix.lower()
    if suffix == ".xlsx":
        return _read_excel_xlsx(path)
    if suffix == ".xls":
        return _read_excel_xls(path)
    return None


def _read_excel_xlsx(path: Path) -> Optional[str]:
    try:
        from openpyxl import load_workbook  # type: ignore[import]
    except ImportError:  # pragma: no cover - dipendenza opzionale
        logger.warning("[filesystem] Missing dependency for xlsx: openpyxl")
        return None

    try:
        workbook = load_workbook(
            filename=str(path),
            data_only=True,
            read_only=True,
        )
    except Exception as exc:
        logger.warning(f"[filesystem] Error reading Excel {path}: {exc}")
        return None

    text_parts: List[str] = []
    try:
        for sheet in workbook.worksheets:
            sheet_lines: List[str] = []
            for row in sheet.iter_rows(values_only=True):
                row_values = [
                    normalize_text(str(value))
                    for value in row
                    if value is not None and str(value).strip()
                ]
                if row_values:
                    sheet_lines.append(" | ".join(row_values))
            if sheet_lines:
                text_parts.append(f"[Foglio {sheet.title}]\n" + "\n".join(sheet_lines))
    finally:
        workbook.close()

    combined = "\n\n".join(text_parts)
    return combined.strip() or None


def _read_excel_xls(path: Path) -> Optional[str]:
    try:
        import xlrd  # type: ignore[import]
    except ImportError:  # pragma: no cover - dipendenza opzionale
        warn_missing_dependency("xlrd", "lettura fogli Excel xls")
        return None

    try:
        workbook = xlrd.open_workbook(str(path))
    except Exception as exc:
        logger.warning(f"[filesystem] Error reading Excel {path}: {exc}")
        return None

    text_parts: List[str] = []
    try:
        for sheet_index in range(workbook.nsheets):
            sheet = workbook.sheet_by_index(sheet_index)
            sheet_lines: List[str] = []
            for row_idx in range(sheet.nrows):
                row_values = [
                    normalize_text(str(sheet.cell_value(row_idx, col_idx)))
                    for col_idx in range(sheet.ncols)
                    if str(sheet.cell_value(row_idx, col_idx)).strip()
                ]
                if row_values:
                    sheet_lines.append(" | ".join(row_values))
            if sheet_lines:
                text_parts.append(f"[Foglio {sheet.name}]\n" + "\n".join(sheet_lines))
    finally:
        try:
            workbook.release_resources()
        except Exception as e:  # pragma: no cover - pulizia best effort
            logger.warning(f"[filesystem] Error releasing XLS resources: {e}")

    combined = "\n\n".join(text_parts)
    return combined.strip() or None
